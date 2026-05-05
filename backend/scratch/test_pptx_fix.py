import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

# Add backend to path so we can import modules
sys.path.append(str(Path(__file__).parent.parent))

from modules.extractor import extract_text_from_pptx, chunk_text

def create_sample_pptx(path):
    prs = Presentation()
    
    # Slide 1: Grouped shapes
    slide1 = prs.slides.add_slide(prs.slide_layouts[6]) # blank
    # Grouped shapes are a bit complex to create programmatically, 
    # so we'll just add multiple shapes and verify they are all picked up.
    # Actually, let's just add a table and notes to slide 1.
    
    # Add a table
    rows, cols = 2, 2
    left, top, width, height = Inches(1), Inches(1), Inches(6), Inches(1.5)
    table = slide1.shapes.add_table(rows, cols, left, top, width, height).table
    table.cell(0, 0).text = "Header 1"
    table.cell(0, 1).text = "Header 2"
    table.cell(1, 0).text = "Value 1"
    table.cell(1, 1).text = "Value 2"
    
    # Add notes
    notes_slide = slide1.notes_slide
    notes_slide.notes_text_frame.text = "This is a secret note."
    
    # Slide 2: Sparse text (to test chunking)
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    tf = txBox.text_frame
    tf.text = "Only ten words in this entire sparse slide today."
    
    prs.save(path)
    print(f"Sample PPTX created at {path}")

def test_extraction(path):
    print("\n--- Testing Extraction ---")
    results = extract_text_from_pptx(path)
    all_text = ""
    for res in results:
        print(f"Page {res['page']}:\n{res['text']}\n")
        all_text += res['text'] + " "
    
    # Assertions
    assert "Header 1" in all_text, "Table header 1 missing"
    assert "Value 2" in all_text, "Table value 2 missing"
    assert "secret note" in all_text, "Notes missing"
    assert "sparse slide" in all_text, "Sparse slide text missing"
    print("Extraction Test Passed!")

def test_chunking():
    print("\n--- Testing Chunking ---")
    sparse_text = "Only ten words in this entire sparse slide today."
    metadata = {"filename": "test.pptx", "page": 1, "file_type": ".pptx"}
    chunks = chunk_text(sparse_text, metadata)
    
    print(f"Chunks found: {len(chunks)}")
    assert len(chunks) > 0, "Chunking failed for sparse text (threshold too high?)"
    assert chunks[0]["text"] == sparse_text, "Chunk text mismatch"
    print("Chunking Test Passed!")

if __name__ == "__main__":
    test_path = "test_sample.pptx"
    try:
        create_sample_pptx(test_path)
        test_extraction(test_path)
        test_chunking()
    finally:
        if os.path.exists(test_path):
            os.remove(test_path)
            print(f"Cleaned up {test_path}")
