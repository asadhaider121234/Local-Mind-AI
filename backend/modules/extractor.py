"""
DocMind — Module 2: Content Extractor
Extracts text from various file formats and chunks it for embedding.
"""

import os
from pathlib import Path
from typing import List, Dict, Any

def dispatch_extractor(file_path: str) -> List[Dict[str, Any]]:
    """
    Route to correct extractor based on extension.
    Return list of: { "page": int, "text": str }
    """
    ext = Path(file_path).suffix.lower()
    
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext in ('.docx', '.doc'):
        return extract_text_from_docx(file_path)
    elif ext in ('.xlsx', '.xls'):
        return extract_text_from_xlsx(file_path)
    elif ext in ('.ppt', '.pptx', '.odp'):
        return extract_text_from_pptx(file_path)
    elif ext == '.csv':
        return extract_text_from_csv(file_path)
    elif ext == '.rtf':
        return extract_text_from_rtf(file_path)
    elif ext in ('.txt', '.md', '.log'):
        return extract_text_from_txt(file_path)
    elif ext in ('.jpg', '.jpeg', '.png'):
        return extract_text_from_image(file_path)
    else:
        # Fallback to code/plain text for everything else
        return extract_text_from_code(file_path)

def extract_text_from_image(file_path: str) -> List[Dict[str, Any]]:
    """
    Combine OCR and Vision Analysis for a comprehensive image index.
    """
    results = []
    
    # 1. OCR for text extraction
    ocr_text = run_ocr_on_file(file_path)
    
    # 2. Vision Analysis (Captioning)
    vision_caption = analyze_image_vision(file_path)
    
    combined_text = ""
    if vision_caption:
        combined_text += f"Image Description: {vision_caption}\n"
    if ocr_text:
        combined_text += f"Extracted Text:\n{ocr_text}"
        
    if combined_text.strip():
        results.append({"page": 1, "text": combined_text.strip()})
    
    return results

def run_ocr_on_file(file_path: str) -> str:
    import pytesseract
    from PIL import Image
    try:
        img = Image.open(file_path)
        return pytesseract.image_to_string(img).strip()
    except Exception as e:
        if "tesseract is not installed" in str(e).lower():
            print(f"[Extractor] OCR error: Tesseract is not installed or not in PATH. Images cannot be indexed.")
        else:
            print(f"[Extractor] OCR error on {file_path}: {e}")
        return ""

def analyze_image_vision(file_path: str) -> str:
    """Uses a lightweight BLIP model to describe the image."""
    vision_type = "BLIP-2 (Recommended)"
    try:
        from api import app_state
        vision_type = app_state.get("vision_model", vision_type)
    except: pass

    if vision_type == "Disabled":
        return ""

    try:
        from PIL import Image
        import torch
        from transformers import BlipProcessor, BlipForConditionalGeneration

        # Note: In a production app, we would cache this model.
        # For the local app, we'll try to use a shared instance if available.
        processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
        model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")

        raw_image = Image.open(file_path).convert('RGB')
        inputs = processor(raw_image, return_tensors="pt")

        out = model.generate(**inputs)
        caption = processor.decode(out[0], skip_special_tokens=True)
        return caption
    except Exception as e:
        print(f"[Extractor] Vision analysis error on {file_path}: {e}")
        return ""

def extract_text_from_pdf(file_path: str) -> List[Dict[str, Any]]:
    import fitz  # PyMuPDF
    results = []
    try:
        doc = fitz.open(file_path)
        for page_num, page in enumerate(doc, start=1):
            text = page.get_text()
            if len(text.strip()) < 50:
                # Fallback to OCR for scanned pages
                pixmap = page.get_pixmap()
                text = run_ocr_on_page(pixmap)
            
            if text.strip():
                results.append({"page": page_num, "text": text.strip()})
        doc.close()
    except Exception as e:
        print(f"[Extractor] Error reading PDF {file_path}: {e}")
    return results

def extract_text_from_docx(file_path: str) -> List[Dict[str, Any]]:
    import docx
    results = []
    try:
        doc = docx.Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        
        # Extract tables as well
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    paragraphs.append(row_text)
                    
        all_text = "\n".join(paragraphs)
        if all_text.strip():
            results.append({"page": 1, "text": all_text.strip()})
    except Exception as e:
        print(f"[Extractor] Error reading DOCX {file_path}: {e}")
    return results

def extract_text_from_xlsx(file_path: str) -> List[Dict[str, Any]]:
    import openpyxl
    results = []
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        all_text_parts = []
        for sheet in wb.worksheets:
            sheet_has_data = False
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    if not sheet_has_data:
                        all_text_parts.append(f"SheetName: {sheet.title}")
                        sheet_has_data = True
                    all_text_parts.append(row_text.strip())
        wb.close()
        
        all_text = "\n".join(all_text_parts)
        if all_text.strip():
            results.append({"page": 1, "text": all_text.strip()})
    except Exception as e:
        print(f"[Extractor] Error reading XLSX {file_path}: {e}")
    return results

def extract_text_from_pptx(file_path: str) -> List[Dict[str, Any]]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    results = []
    
    # Check for legacy .ppt format
    if file_path.lower().endswith('.ppt'):
        print(f"[Extractor] Warning: Legacy .ppt format detected for {file_path}. Only .pptx is supported by the current engine.")
        return results

    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides, start=1):
            text_parts = []
            
            def extract_from_shape(shape):
                parts = []
                # 1. Standard text
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
                
                # 2. Tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text_frame.text.strip() for cell in row.cells if cell.text_frame.text.strip())
                        if row_text:
                            parts.append(row_text)
                
                # 3. Recursive Groups
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for subshape in shape.shapes:
                        parts.extend(extract_from_shape(subshape))
                return parts

            for shape in slide.shapes:
                text_parts.extend(extract_from_shape(shape))
            
            # 4. Slide Notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    text_parts.append(f"Notes: {notes}")
            
            slide_text = "\n".join(text_parts)
            if slide_text.strip():
                results.append({"page": i, "text": slide_text.strip()})
    except Exception as e:
        print(f"[Extractor] Error reading PowerPoint {file_path}: {e}")
    return results

def extract_text_from_rtf(file_path: str) -> List[Dict[str, Any]]:
    from striprtf.striprtf import rtf_to_text
    results = []
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            rtf_content = f.read()
        text = rtf_to_text(rtf_content)
        if text.strip():
            results.append({"page": 1, "text": text.strip()})
    except Exception as e:
        print(f"[Extractor] Error reading RTF {file_path}: {e}")
    return results

def extract_text_from_txt(file_path: str) -> List[Dict[str, Any]]:
    results = []
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            all_text = f.read()
    except UnicodeDecodeError:
        try:
            with open(file_path, "r", encoding="latin-1") as f:
                all_text = f.read()
        except Exception as e:
            print(f"[Extractor] Error reading TXT {file_path}: {e}")
            return results
    except Exception as e:
        print(f"[Extractor] Error reading TXT {file_path}: {e}")
        return results

    if all_text.strip():
        results.append({"page": 1, "text": all_text.strip()})
    return results

def extract_text_from_csv(file_path: str) -> List[Dict[str, Any]]:
    import pandas as pd
    results = []
    try:
        df = pd.read_csv(file_path, nrows=1000)
        # Convert df to string representation
        all_text = df.to_string(index=False)
        if all_text.strip():
            results.append({"page": 1, "text": all_text.strip()})
    except Exception as e:
        print(f"[Extractor] Error reading CSV {file_path}: {e}")
    return results

def extract_text_from_code(file_path: str) -> List[Dict[str, Any]]:
    results = []
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            all_text = f.read()
        if all_text.strip():
            results.append({"page": 1, "text": all_text.strip()})
    except Exception as e:
        print(f"[Extractor] Error reading code/text {file_path}: {e}")
    return results

def run_ocr_on_page(pixmap) -> str:
    import pytesseract
    from PIL import Image
    import io
    try:
        img_data = pixmap.tobytes("png")
        img = Image.open(io.BytesIO(img_data))
        text = pytesseract.image_to_string(img)
        return text
    except Exception as e:
        if "tesseract is not installed" in str(e).lower():
            print(f"[Extractor] OCR error: Tesseract missing. Scanned PDF pages will be skipped.")
        else:
            print(f"[Extractor] OCR error: {e}")
        return ""

def chunk_text(text: str, metadata: dict, chunk_size=400, overlap=80) -> List[Dict[str, Any]]:
    """
    Split text into overlapping word-based chunks.
    metadata should contain: filename, page, file_type
    """
    chunks = []
    words = text.split()
    
    if not words:
        return chunks
        
    start_idx = 0
    chunk_id = 0
    
    # Reduced threshold from 30 to 5 to avoid losing sparse slides/images
    MIN_WORDS = 5
    
    while start_idx < len(words):
        end_idx = min(start_idx + chunk_size, len(words))
        chunk_words = words[start_idx:end_idx]
        
        if len(chunk_words) >= MIN_WORDS:
            chunk_text_str = " ".join(chunk_words)
            chunks.append({
                "text": chunk_text_str,
                "filename": metadata.get("filename", ""),
                "page": metadata.get("page", 1),
                "file_type": metadata.get("file_type", ""),
                "category": "",
                "chunk_id": chunk_id
            })
            chunk_id += 1
            
        if end_idx == len(words):
            break
            
        start_idx = end_idx - overlap

    return chunks
